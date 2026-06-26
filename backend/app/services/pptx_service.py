"""PowerPoint <-> PDF conversion service.

- powerpoint_to_pdf: uses LibreOffice headless (same binary as word-to-pdf)
- pdf_to_powerpoint: renders each PDF page to an image via pdf2image,
  then drops each image as a full-bleed picture onto its own slide via
  python-pptx. Text/vector fidelity is not preserved (the result is a
  slide-deck of page screenshots) — matches what iLovePDF does, and is
  what users actually want when they say "convert PDF to PowerPoint".
"""
from __future__ import annotations

import io
import os
import subprocess
import tempfile
from pathlib import Path

from app.utils import PDFProcessingError
from app.utils.concurrency import resolve_libreoffice_path


class PptxService:
    @staticmethod
    def pptx_to_pdf(pptx_bytes: bytes, original_filename: str) -> bytes:
        return powerpoint_to_pdf(pptx_bytes, original_filename)

    @staticmethod
    def pdf_to_pptx(pdf_bytes: bytes, dpi: int = 150) -> bytes:
        return pdf_to_powerpoint(pdf_bytes, dpi=dpi)


def powerpoint_to_pdf(pptx_bytes: bytes, original_filename: str) -> bytes:
    """Convert a .ppt/.pptx file to PDF using LibreOffice headless."""
    libreoffice_path = resolve_libreoffice_path()
    if not libreoffice_path:
        raise PDFProcessingError("LibreOffice is not installed on the server.")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        input_path = tmp_path / original_filename
        input_path.write_bytes(pptx_bytes)

        # Unique profile dir so parallel LibreOffice runs don't collide.
        import uuid as _uuid
        profile_dir = tmp_path / f"lo_profile_{_uuid.uuid4().hex[:8]}"
        cmd = [
            libreoffice_path,
            "--headless",
            f"-env:UserInstallation=file:///{profile_dir.as_posix().lstrip('/')}",
            "--convert-to", "pdf",
            "--outdir", str(tmp_path),
            str(input_path),
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        except subprocess.TimeoutExpired:
            raise PDFProcessingError("PowerPoint to PDF conversion timed out.")

        if result.returncode != 0:
            raise PDFProcessingError(
                f"LibreOffice failed: {result.stderr[:200] or 'unknown error'}"
            )

        pdfs = list(tmp_path.glob("*.pdf"))
        if not pdfs:
            raise PDFProcessingError("Conversion completed but no PDF was generated.")
        return pdfs[0].read_bytes()


def pdf_to_powerpoint(pdf_bytes: bytes, dpi: int = 150) -> bytes:
    """Convert a PDF to a .pptx where each page becomes one slide
    containing a full-bleed screenshot of that page.
    """
    try:
        from pdf2image import convert_from_bytes
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError as e:
        raise PDFProcessingError(
            f"Missing dependency for PDF to PowerPoint: {e}. "
            "Install python-pptx and pdf2image."
        )

    # Render PDF pages to PIL images. POPPLER_PATH is set in app settings on Windows.
    poppler_path = os.environ.get("POPPLER_PATH") or r"C:/shahrukh/poppler/Library/bin"
    try:
        images = convert_from_bytes(pdf_bytes, dpi=dpi, poppler_path=poppler_path)
    except Exception as e:
        raise PDFProcessingError(f"Failed to render PDF pages: {e}")

    if not images:
        raise PDFProcessingError("PDF has no pages to convert.")

    # 16:9 slide is the modern default — 13.333" x 7.5"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    blank_layout = prs.slide_layouts[6]  # blank

    for img in images:
        slide = prs.slides.add_slide(blank_layout)

        # Fit image into the slide while preserving aspect ratio (letterboxed).
        img_w, img_h = img.size
        slide_ratio = slide_w_emu / slide_h_emu
        img_ratio = img_w / img_h
        if img_ratio > slide_ratio:
            # Image is wider than slide -> fit by width
            new_w = slide_w_emu
            new_h = Emu(int(slide_w_emu * img_h / img_w))
            left = Emu(0)
            top = Emu(int((slide_h_emu - new_h) / 2))
        else:
            # Image is taller than slide -> fit by height
            new_h = slide_h_emu
            new_w = Emu(int(slide_h_emu * img_w / img_h))
            left = Emu(int((slide_w_emu - new_w) / 2))
            top = Emu(0)

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, left, top, width=new_w, height=new_h)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
